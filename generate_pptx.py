from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette (matches landing page zinc/dark theme) ──────────────────────────
BG        = RGBColor(0x09, 0x09, 0x0B)   # near-black  #09090B
CARD      = RGBColor(0x18, 0x18, 0x1B)   # zinc-900    #18181B
CARD2     = RGBColor(0x27, 0x27, 0x2A)   # zinc-800    #27272A
BORDER    = RGBColor(0x3F, 0x3F, 0x46)   # zinc-700    #3F3F46
FG        = RGBColor(0xFF, 0xFF, 0xFF)   # white
MUTED     = RGBColor(0xA1, 0xA1, 0xAA)   # zinc-400    #A1A1AA
MUTED2    = RGBColor(0x71, 0x71, 0x7A)   # zinc-500    #71717A
GREEN     = RGBColor(0x4A, 0xDE, 0x80)   # green-400   #4ADE80
GREEN_BG  = RGBColor(0x05, 0x2E, 0x16)   # green-950
GRADIENT1 = RGBColor(0xE4, 0xE4, 0xE7)   # zinc-200
GRADIENT2 = RGBColor(0x71, 0x71, 0x7A)   # zinc-500

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    bg_fill(slide)
    return slide

def bg_fill(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def box(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0.75), corner=0.08):
    shape = slide.shapes.add_shape(
        5,  # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        x, y, w, h
    )
    shape.adjustments[0] = corner
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=18, bold=False, color=FG, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    tf.margin_left   = 0
    tf.margin_right  = 0
    tf.margin_top    = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = "Segoe UI"
    return txb

def bullet_row(slide, text, x, y, dot_color=MUTED2):
    dot = slide.shapes.add_shape(1, x, y + Inches(0.08), Inches(0.07), Inches(0.07))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()
    txt(slide, text, x + Inches(0.15), y, Inches(3.8), Inches(0.3), size=11, color=MUTED)

def divider(slide, y, left=Inches(0.5), right=Inches(12.83)):
    line = slide.shapes.add_connector(1, left, y, right, y)
    line.line.color.rgb = BORDER
    line.line.width = Pt(0.5)

def pill(slide, text, x, y, bg=CARD2, fg=MUTED):
    shape = box(slide, x, y, Inches(1.5), Inches(0.28), fill_color=bg, line_color=BORDER, line_width=Pt(0.5), corner=0.5)
    txt(slide, text, x + Inches(0.1), y + Inches(0.03), Inches(1.3), Inches(0.22), size=9, color=fg, align=PP_ALIGN.CENTER)
    return shape

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()

# Green pulse badge
badge_bg = box(s, Inches(4.9), Inches(1.1), Inches(3.5), Inches(0.38), fill_color=CARD2, line_color=BORDER, line_width=Pt(0.5))
dot = s.shapes.add_shape(1, Inches(5.12), Inches(1.24), Inches(0.09), Inches(0.09))
dot.fill.solid(); dot.fill.fore_color.rgb = GREEN; dot.line.fill.background()
txt(s, "Multi-Tenant SaaS Platform", Inches(5.25), Inches(1.12), Inches(3.1), Inches(0.35),
    size=11, color=RGBColor(0xD4,0xD4,0xD8), align=PP_ALIGN.LEFT)

# Main title
txt(s, "Transform your", Inches(2.0), Inches(1.65), Inches(9.3), Inches(0.85),
    size=52, bold=True, color=FG, align=PP_ALIGN.CENTER)
txt(s, "motorcycle shop", Inches(2.0), Inches(2.45), Inches(9.3), Inches(0.85),
    size=52, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, "operations", Inches(2.0), Inches(3.25), Inches(9.3), Inches(0.85),
    size=52, bold=True, color=FG, align=PP_ALIGN.CENTER)

# Subtitle
txt(s, "All-in-one SaaS platform for inventory, service jobs, sales, reports, and team management.\nMulti-tenant by design — each shop gets its own branded subdomain.",
    Inches(2.8), Inches(4.2), Inches(7.7), Inches(0.8),
    size=14, color=MUTED, align=PP_ALIGN.CENTER)

# Trust pills
for i, (ico, label) in enumerate([("⚡", "Secure & Reliable"), ("⚡", "Real-time Sync"), ("↑", "Analytics Included")]):
    px = Inches(3.5) + i * Inches(2.2)
    pill(s, label, px, Inches(5.3), bg=CARD2, fg=MUTED)

# Product name bottom-left watermark
txt(s, "MoSPAMS", Inches(0.3), Inches(7.0), Inches(2.0), Inches(0.35),
    size=10, color=MUTED2, bold=False)
txt(s, "mospams.shop", Inches(10.8), Inches(7.0), Inches(2.2), Inches(0.35),
    size=10, color=MUTED2, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — About MoSPAMS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()

# Section label
pill(s, "About MoSPAMS", Inches(0.5), Inches(0.45), bg=CARD2, fg=MUTED)

txt(s, "Built for small-to-medium", Inches(0.5), Inches(0.95), Inches(7.0), Inches(0.65),
    size=34, bold=True, color=FG)
txt(s, "motorcycle businesses", Inches(0.5), Inches(1.55), Inches(7.0), Inches(0.65),
    size=34, bold=True, color=MUTED)

txt(s, "MoSPAMS is designed to reduce manual paperwork, prevent stock\ndiscrepancies, organize service records, speed up transactions, and\nprovide useful business insights for motorcycle repair shops and parts retailers.",
    Inches(0.5), Inches(2.35), Inches(6.2), Inches(1.0),
    size=13, color=MUTED)

# Benefit cards (left column)
benefits = [
    ("Reduce Manual Errors",
     "Eliminate handwritten records and spreadsheet mistakes. MoSPAMS keeps all data accurate."),
    ("Improve Shop Workflow",
     "Speed up service jobs, transactions, and stock updates so your team does more real work."),
    ("Make Better Business Decisions",
     "Use sales reports, inventory summaries, and service data to grow your shop smarter."),
]
for i, (title, desc) in enumerate(benefits):
    by = Inches(3.5) + i * Inches(1.1)
    card = box(s, Inches(0.5), by, Inches(6.0), Inches(0.95), fill_color=CARD, line_color=BORDER, line_width=Pt(0.5))
    txt(s, title, Inches(0.7), by + Inches(0.08), Inches(5.7), Inches(0.28), size=12, bold=True, color=FG)
    txt(s, desc,  Inches(0.7), by + Inches(0.35), Inches(5.7), Inches(0.5),  size=10, color=MUTED)

# Stats (right column)
stats = [("100%", "Cloud-Based"), ("4", "User Roles"), ("5+", "Report Types"), ("∞", "Transactions")]
for i, (val, lbl) in enumerate(stats):
    col = i % 2
    row = i // 2
    sx = Inches(7.2) + col * Inches(2.8)
    sy = Inches(0.9)  + row * Inches(1.5)
    card = box(s, sx, sy, Inches(2.5), Inches(1.2), fill_color=CARD, line_color=BORDER, line_width=Pt(0.5))
    txt(s, val, sx, sy + Inches(0.15), Inches(2.5), Inches(0.6), size=36, bold=True, color=FG, align=PP_ALIGN.CENTER)
    txt(s, lbl, sx, sy + Inches(0.72), Inches(2.5), Inches(0.35), size=12, color=MUTED, align=PP_ALIGN.CENTER)

# Version card
vc = box(s, Inches(7.2), Inches(3.95), Inches(5.6), Inches(2.9), fill_color=CARD, line_color=BORDER, line_width=Pt(0.5))
txt(s, "Version 1.0", Inches(7.5), Inches(4.1), Inches(5.0), Inches(0.4), size=18, bold=True, color=FG)
txt(s, "SaaS Platform", Inches(7.5), Inches(4.5), Inches(5.0), Inches(0.3), size=12, color=MUTED)
txt(s, "MoSPAMS is a full multi-tenant SaaS platform. Owner, Staff,\nMechanic, and Customer roles are all active with dedicated\ndashboards and workflows.",
    Inches(7.5), Inches(4.85), Inches(5.0), Inches(0.8), size=11, color=MUTED)
tags = ["Inventory","Services","Sales","Reports","Users","Branding","Multi-Tenancy","Google Auth"]
for i, t in enumerate(tags):
    tx = Inches(7.5) + (i % 4) * Inches(1.3)
    ty = Inches(5.75) + (i // 4) * Inches(0.32)
    pill(s, t, tx, ty, bg=CARD2, fg=MUTED)

txt(s, "MoSPAMS", Inches(0.3), Inches(7.0), Inches(2.0), Inches(0.35), size=10, color=MUTED2)
txt(s, "mospams.shop", Inches(10.8), Inches(7.0), Inches(2.2), Inches(0.35), size=10, color=MUTED2, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Core Features
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()

pill(s, "Core Features", Inches(0.5), Inches(0.4), bg=CARD2, fg=MUTED)
txt(s, "Powerful features built for motorcycle shop management",
    Inches(0.5), Inches(0.88), Inches(12.3), Inches(0.5),
    size=26, bold=True, color=FG, align=PP_ALIGN.CENTER)
txt(s, "Everything your shop needs — inventory, service jobs, team management, and branding.",
    Inches(1.5), Inches(1.42), Inches(10.3), Inches(0.35),
    size=13, color=MUTED, align=PP_ALIGN.CENTER)

features = [
    ("Inventory Management",
     "Track parts, categories, stock movements, low-stock alerts, and barcode lookups.",
     ["Stock Tracking", "Categories", "Low-Stock Alerts"]),
    ("Service Job Tracking",
     "Create service records, assign mechanics, attach parts, and track job status.",
     ["Mechanic Assignment", "Job Parts", "Status Flow"]),
    ("Sales & Transactions",
     "Record parts-only and service transactions with Cash and GCash payment tracking.",
     ["Cash", "GCash", "Discounts"]),
    ("Reports & Analytics",
     "Sales reports, inventory summaries, service performance, and real-time KPI dashboard.",
     ["Sales Report", "Income", "Dashboard KPIs"]),
    ("Role-Based Access",
     "Five distinct roles — Owner, Staff, Mechanic, and Customer — each with tailored permissions.",
     ["5 Roles", "Permissions", "Google Sign-In"]),
    ("Multi-Tenant & Branding",
     "Each shop gets its own subdomain, logo, color scheme, and fully isolated data.",
     ["Subdomains", "Shop Branding", "Data Isolation"]),
]

cols = 3
for idx, (title, desc, tags) in enumerate(features):
    col = idx % cols
    row = idx // cols
    fx = Inches(0.35) + col * Inches(4.3)
    fy = Inches(2.0)  + row * Inches(2.55)
    card = box(s, fx, fy, Inches(4.1), Inches(2.35), fill_color=CARD, line_color=BORDER, line_width=Pt(0.5))
    # icon placeholder dot
    dot = s.shapes.add_shape(1, fx + Inches(0.18), fy + Inches(0.18), Inches(0.3), Inches(0.3))
    dot.fill.solid(); dot.fill.fore_color.rgb = CARD2; dot.line.color.rgb = BORDER; dot.line.width = Pt(0.5)
    txt(s, title, fx + Inches(0.18), fy + Inches(0.58), Inches(3.7), Inches(0.32), size=13, bold=True, color=FG)
    txt(s, desc,  fx + Inches(0.18), fy + Inches(0.92), Inches(3.7), Inches(0.65), size=10, color=MUTED)
    # Tags: distribute evenly across inner card width (3.74") with small gaps
    tag_count = len(tags)
    tag_gap   = Inches(0.07)
    tag_w     = (Inches(3.74) - tag_gap * (tag_count - 1)) / tag_count
    for ti, tag in enumerate(tags):
        tx = fx + Inches(0.18) + ti * (tag_w + tag_gap)
        ty = fy + Inches(1.9)
        box(s, tx, ty, tag_w, Inches(0.28), fill_color=CARD2, line_color=BORDER, line_width=Pt(0.5), corner=0.5)
        txt(s, tag, tx + Inches(0.07), ty + Inches(0.03), tag_w - Inches(0.14), Inches(0.22), size=8, color=MUTED, align=PP_ALIGN.CENTER)

txt(s, "MoSPAMS", Inches(0.3), Inches(7.0), Inches(2.0), Inches(0.35), size=10, color=MUTED2)
txt(s, "mospams.shop", Inches(10.8), Inches(7.0), Inches(2.2), Inches(0.35), size=10, color=MUTED2, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Live Dashboard Overview
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()

pill(s, "Dashboard Preview", Inches(0.5), Inches(0.4), bg=CARD2, fg=MUTED)
txt(s, "Real-time shop overview at a glance",
    Inches(0.5), Inches(0.88), Inches(8.0), Inches(0.5),
    size=26, bold=True, color=FG)
txt(s, "KPI cards, sales charts, and live data — all in one unified dashboard.",
    Inches(0.5), Inches(1.42), Inches(8.0), Inches(0.35),
    size=13, color=MUTED)

# Dashboard card container
db = box(s, Inches(0.4), Inches(1.95), Inches(12.5), Inches(5.1), fill_color=CARD, line_color=BORDER, line_width=Pt(0.75))

# Header bar inside dashboard
txt(s, "MoSPAMS Dashboard", Inches(0.7), Inches(2.05), Inches(4.0), Inches(0.3), size=12, bold=True, color=FG)
txt(s, "Real-time overview", Inches(0.7), Inches(2.35), Inches(4.0), Inches(0.25), size=10, color=MUTED)
# Live badge
live_bg = box(s, Inches(11.0), Inches(2.1), Inches(1.6), Inches(0.28), fill_color=GREEN_BG, line_color=RGBColor(0x14,0x53,0x2D), line_width=Pt(0.5))
dot2 = s.shapes.add_shape(1, Inches(11.18), Inches(2.18), Inches(0.1), Inches(0.1))
dot2.fill.solid(); dot2.fill.fore_color.rgb = GREEN; dot2.line.fill.background()
txt(s, "Live", Inches(11.32), Inches(2.1), Inches(0.8), Inches(0.28), size=10, color=GREEN)

divider(s, Inches(2.5), left=Inches(0.5), right=Inches(12.8))

# KPI stat cards
kpis = [("Revenue", "PHP 48.2K", "+12%"), ("Active Jobs", "14", "+3"), ("Parts Stock", "1,247", "-8"), ("Completed", "128", "+8")]
for i, (label, val, trend) in enumerate(kpis):
    kx = Inches(0.65) + i * Inches(3.05)
    ky = Inches(2.65)
    kcard = box(s, kx, ky, Inches(2.8), Inches(1.0), fill_color=CARD2, line_color=BORDER, line_width=Pt(0.5))
    txt(s, label, kx + Inches(0.15), ky + Inches(0.07), Inches(2.5), Inches(0.24), size=10, color=MUTED)
    txt(s, val,   kx + Inches(0.15), ky + Inches(0.33), Inches(2.5), Inches(0.35), size=20, bold=True, color=FG)
    tcolor = GREEN if trend.startswith("+") else MUTED2
    txt(s, trend, kx + Inches(0.15), ky + Inches(0.68), Inches(2.5), Inches(0.24), size=10, bold=True, color=tcolor)

# Sales chart area
chart_bg = box(s, Inches(0.65), Inches(3.82), Inches(11.95), Inches(2.85), fill_color=CARD2, line_color=BORDER, line_width=Pt(0.5))
txt(s, "Sales Performance", Inches(0.9), Inches(3.95), Inches(4.0), Inches(0.28), size=10, color=MUTED)
txt(s, "PHP 125,450",       Inches(0.9), Inches(4.25), Inches(4.0), Inches(0.4),  size=22, bold=True, color=FG)
# Growth badge
gb = box(s, Inches(11.1), Inches(3.95), Inches(1.2), Inches(0.28), fill_color=GREEN_BG, line_color=RGBColor(0x14,0x53,0x2D), line_width=Pt(0.5))
txt(s, "+18.2%", Inches(11.15), Inches(3.95), Inches(1.1), Inches(0.28), size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Bar chart
heights = [0.45, 0.70, 0.55, 0.85, 0.60, 0.95, 0.75, 0.90, 0.65, 0.88, 0.70, 0.92, 0.78, 0.95, 0.82, 0.88]
chart_max_h = Inches(1.4)
chart_base  = Inches(6.45)
for i, h in enumerate(heights):
    bx = Inches(0.85) + i * Inches(0.74)
    bh = chart_max_h * h
    by = chart_base - bh
    bar = s.shapes.add_shape(1, bx, by, Inches(0.55), bh)
    bar.fill.solid(); bar.fill.fore_color.rgb = BORDER; bar.line.fill.background()

for i, lbl in enumerate(["Jan", "Feb", "Mar", "Apr"]):
    txt(s, lbl, Inches(1.5) + i * Inches(2.95), Inches(6.52), Inches(1.0), Inches(0.25), size=9, color=MUTED2)

txt(s, "MoSPAMS", Inches(0.3), Inches(7.0), Inches(2.0), Inches(0.35), size=10, color=MUTED2)
txt(s, "mospams.shop", Inches(10.8), Inches(7.0), Inches(2.2), Inches(0.35), size=10, color=MUTED2, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — User Roles
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()

pill(s, "User Roles", Inches(0.5), Inches(0.4), bg=CARD2, fg=MUTED)
txt(s, "Designed for real shop workflows",
    Inches(0.5), Inches(0.88), Inches(12.3), Inches(0.5),
    size=26, bold=True, color=FG, align=PP_ALIGN.CENTER)
txt(s, "Role-based access ensures every team member sees exactly what they need — nothing more, nothing less.",
    Inches(1.0), Inches(1.42), Inches(11.3), Inches(0.35),
    size=13, color=MUTED, align=PP_ALIGN.CENTER)

roles = [
    ("Owner", "Primary", "Full shop access: inventory, service jobs, sales, reports, user management, shop branding, and activity logs.",
     ["Inventory Management", "Service Jobs", "Sales & Transactions", "Reports & Analytics", "User Management", "Shop Branding", "Activity Logs"]),
    ("Staff", "Active", "Operational access to daily shop tasks: manage inventory, service jobs, sales transactions, and view reports.",
     ["Inventory Management", "Stock Movements", "Service Job Mgmt", "Sales & Transactions", "View Reports"]),
    ("Mechanic", "Active", "Dedicated dashboard for assigned service jobs. View job details, update status, add or remove parts used.",
     ["View Assigned Jobs", "Update Job Status", "Add/Remove Parts", "Job Details"]),
    ("Customer", "Active", "Self-service portal: view service history, submit service requests, and track payment records.",
     ["View Service History", "Create Service Requests", "Payment History"]),
]

for idx, (role, badge, desc, perms) in enumerate(roles):
    rx = Inches(0.35) + idx * Inches(3.2)
    ry = Inches(2.0)
    rcard = box(s, rx, ry, Inches(3.0), Inches(4.85), fill_color=CARD, line_color=BORDER if badge != "Primary" else RGBColor(0x52,0x52,0x5B), line_width=Pt(0.75))
    # card header
    hdr = box(s, rx, ry, Inches(3.0), Inches(1.0), fill_color=CARD2, line_color=None)
    txt(s, role, rx + Inches(0.18), ry + Inches(0.15), Inches(2.6), Inches(0.38), size=16, bold=True, color=FG)
    # badge
    bg_b = GREEN_BG if badge == "Active" or badge == "Primary" else CARD2
    fg_b = GREEN if badge == "Active" or badge == "Primary" else MUTED
    pill(s, badge, rx + Inches(0.18), ry + Inches(0.6), bg=bg_b, fg=fg_b)
    if badge == "Primary":
        pill(s, "Primary", rx + Inches(1.82), ry + Inches(0.04), bg=FG, fg=BG)

    txt(s, desc, rx + Inches(0.18), ry + Inches(1.1), Inches(2.65), Inches(0.75), size=10, color=MUTED)
    for pi, perm in enumerate(perms):
        py = ry + Inches(1.95) + pi * Inches(0.38)
        dot3 = s.shapes.add_shape(1, rx + Inches(0.18), py + Inches(0.1), Inches(0.09), Inches(0.09))
        dot3.fill.solid(); dot3.fill.fore_color.rgb = GREEN; dot3.line.fill.background()
        txt(s, perm, rx + Inches(0.35), py, Inches(2.5), Inches(0.3), size=10, color=RGBColor(0xD4,0xD4,0xD8))

# Note banner
nb = box(s, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.6), fill_color=CARD, line_color=BORDER, line_width=Pt(0.5))
txt(s, "Role-Based Onboarding:  New team members sign up with Google and request their role. The Owner approves each request before access is granted.",
    Inches(0.65), Inches(6.72), Inches(12.0), Inches(0.45), size=10, color=MUTED)

txt(s, "MoSPAMS", Inches(0.3), Inches(7.1), Inches(2.0), Inches(0.35), size=10, color=MUTED2)
txt(s, "mospams.shop", Inches(10.8), Inches(7.1), Inches(2.2), Inches(0.35), size=10, color=MUTED2, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Why MoSPAMS? (Benefits summary)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()

pill(s, "Why MoSPAMS?", Inches(0.5), Inches(0.4), bg=CARD2, fg=MUTED)
txt(s, "The smartest way to run your motorcycle shop",
    Inches(0.5), Inches(0.88), Inches(12.3), Inches(0.5),
    size=28, bold=True, color=FG, align=PP_ALIGN.CENTER)

benefits2 = [
    ("Stop Losing Money to Manual Errors",
     "Handwritten records and spreadsheets cause stock discrepancies, missed charges, and lost revenue. MoSPAMS eliminates these by keeping every part, job, and transaction in one accurate system.",
     ["Automated stock updates", "Accurate invoices", "Zero double-entry"]),
    ("Run a Faster, More Organized Shop",
     "From creating a service job to processing payment, MoSPAMS shortens every workflow. Mechanics see their jobs instantly, staff record sales in seconds, and owners get reports with one click.",
     ["Mechanic job queue", "Fast POS-style sales", "One-click reports"]),
    ("Grow With Real Data",
     "MoSPAMS gives you the business intelligence to make smart decisions — which parts sell most, which services earn most, and which mechanics are most productive.",
     ["Sales trend reports", "Income breakdowns", "Performance KPIs"]),
]

for i, (title, body, pts) in enumerate(benefits2):
    bx = Inches(0.4)
    by = Inches(1.75) + i * Inches(1.7)
    card = box(s, bx, by, Inches(12.5), Inches(1.5), fill_color=CARD, line_color=BORDER, line_width=Pt(0.5))
    num = box(s, bx + Inches(0.15), by + Inches(0.2), Inches(0.55), Inches(0.55), fill_color=CARD2, line_color=BORDER, line_width=Pt(0.5))
    txt(s, str(i+1), bx + Inches(0.15), by + Inches(0.2), Inches(0.55), Inches(0.55), size=20, bold=True, color=FG, align=PP_ALIGN.CENTER)
    txt(s, title, bx + Inches(0.85), by + Inches(0.1), Inches(5.5), Inches(0.38), size=14, bold=True, color=FG)
    txt(s, body,  bx + Inches(0.85), by + Inches(0.5), Inches(5.5), Inches(0.8), size=10, color=MUTED)
    for pi, pt in enumerate(pts):
        ppx = Inches(7.2) + pi * Inches(1.75)
        ppy = by + Inches(0.55)
        pill(s, pt, ppx, ppy, bg=CARD2, fg=MUTED)

# Bottom CTA teaser
cta_bg = box(s, Inches(1.5), Inches(6.55), Inches(10.3), Inches(0.65), fill_color=CARD2, line_color=BORDER, line_width=Pt(0.5))
txt(s, "Ready to transform your shop?  →  mospams.shop",
    Inches(1.5), Inches(6.62), Inches(10.3), Inches(0.5),
    size=14, bold=True, color=FG, align=PP_ALIGN.CENTER)

txt(s, "MoSPAMS", Inches(0.3), Inches(7.1), Inches(2.0), Inches(0.35), size=10, color=MUTED2)
txt(s, "mospams.shop", Inches(10.8), Inches(7.1), Inches(2.2), Inches(0.35), size=10, color=MUTED2, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Call to Action / Thank You
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()

txt(s, "Get started today —", Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.75),
    size=46, bold=True, color=FG, align=PP_ALIGN.CENTER)
txt(s, "your shop deserves better tools.", Inches(1.5), Inches(2.55), Inches(10.3), Inches(0.75),
    size=46, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

txt(s, "MoSPAMS — Motorcycle Service and Parts Management System",
    Inches(2.0), Inches(3.55), Inches(9.3), Inches(0.4),
    size=14, color=MUTED, align=PP_ALIGN.CENTER)

# CTA buttons (visual only)
btn1 = box(s, Inches(3.8), Inches(4.2), Inches(2.6), Inches(0.5), fill_color=FG, line_color=None, corner=0.2)
txt(s, "Start free trial", Inches(3.8), Inches(4.25), Inches(2.6), Inches(0.4),
    size=13, bold=True, color=BG, align=PP_ALIGN.CENTER)

btn2 = box(s, Inches(6.9), Inches(4.2), Inches(2.6), Inches(0.5), fill_color=CARD2, line_color=BORDER, line_width=Pt(0.5), corner=0.2)
txt(s, "View Features", Inches(6.9), Inches(4.25), Inches(2.6), Inches(0.4),
    size=13, bold=True, color=FG, align=PP_ALIGN.CENTER)

txt(s, "mospams.shop", Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.45),
    size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# Contact row
txt(s, "support@mospams.shop  •  14-day free trial  •  No credit card required",
    Inches(2.0), Inches(5.85), Inches(9.3), Inches(0.35),
    size=12, color=MUTED2, align=PP_ALIGN.CENTER)

# Decorative bottom line
divider(s, Inches(6.55), left=Inches(2.0), right=Inches(11.3))

txt(s, "MoSPAMS — Confidential Product Overview", Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3),
    size=10, color=MUTED2, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out = r"C:\Users\frien\Documents\Rizal_V1 MVP\MoSPAMS\MoSPAMS_Product_Overview.pptx"
prs.save(out)
print(f"Saved: {out}")
